#!/usr/bin/env python3
"""Stitch에서 내보낸 슬라이드 이미지를 16:9 pptx로 조립한다.

사용법:
    python assemble_pptx.py <이미지폴더> -o 발표자료.pptx
    python assemble_pptx.py <이미지폴더> -o 발표자료.pptx --notes notes.txt
    python assemble_pptx.py <이미지폴더> -o 발표자료.pptx --overlay overlay.json

- 이미지 폴더 안의 png/jpg/jpeg/webp 파일을 자연 정렬 순서로 슬라이드에 넣는다.
- 슬라이드 크기는 16:9 (13.333 x 7.5 인치).
- 이미지 비율이 16:9와 다르면 여백 없이 중앙을 채우도록(cover) 크롭 배치한다.
- --notes 파일은 슬라이드별 발표자 노트. 슬라이드 사이를 '---' 한 줄로 구분한다.
- --overlay 파일은 사용자 첨부 이미지를 슬라이드 위에 얹는 배치표(JSON).
  자세한 규격은 references/image-slides.md 참고.
"""
import argparse
import json
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit("python-pptx가 없습니다. 먼저 실행: pip install python-pptx --break-system-packages")

try:
    from PIL import Image
except ImportError:
    sys.exit("Pillow가 없습니다. 먼저 실행: pip install Pillow --break-system-packages")

SLIDE_W = Emu(12192000)  # 13.333 inch
SLIDE_H = Emu(6858000)   # 7.5 inch
EXTS = {".png", ".jpg", ".jpeg", ".webp"}

# 배치 영역 프리셋: (x, y, w, h) — 슬라이드 크기 대비 0~1 비율
BOXES = {
    "full":         (0.0,  0.0,  1.0,  1.0),
    "left-half":    (0.0,  0.0,  0.5,  1.0),
    "right-half":   (0.5,  0.0,  0.5,  1.0),
    "top-half":     (0.0,  0.0,  1.0,  0.5),
    "bottom-half":  (0.0,  0.5,  1.0,  0.5),
    "top-72":       (0.0,  0.0,  1.0,  0.72),   # image-full 기본값(하단 28%는 텍스트 밴드)
    "left-third":   (0.0,  0.0,  0.34, 1.0),
    "right-third":  (0.66, 0.0,  0.34, 1.0),
    "left-panel":   (0.06, 0.28, 0.42, 0.50),   # image-compare 왼쪽
    "right-panel":  (0.52, 0.28, 0.42, 0.50),   # image-compare 오른쪽
    "center":       (0.20, 0.28, 0.60, 0.52),   # logo-grid 그리드 영역
}


def natural_key(p: Path):
    return [int(t) if t.isdigit() else t.lower() for t in re.split(r"(\d+)", p.name)]


def load_notes(path: Path):
    text = path.read_text(encoding="utf-8")
    return [chunk.strip() for chunk in re.split(r"^\s*---\s*$", text, flags=re.M)]


def resolve_box(spec):
    """프리셋 이름 또는 [x, y, w, h] 비율을 EMU 사각형으로 바꾼다."""
    if spec is None:
        spec = "full"
    if isinstance(spec, str):
        if spec not in BOXES:
            sys.exit(f"알 수 없는 box 이름: {spec}\n사용 가능: {', '.join(sorted(BOXES))}")
        ratio = BOXES[spec]
    else:
        if len(spec) != 4:
            sys.exit(f"box 좌표는 [x, y, w, h] 네 개여야 합니다: {spec}")
        ratio = tuple(float(v) for v in spec)
    x, y, w, h = ratio
    return (int(SLIDE_W * x), int(SLIDE_H * y), int(SLIDE_W * w), int(SLIDE_H * h))


def place_image(slide, img_path: Path, box, fit: str):
    """box(EMU 사각형) 안에 이미지를 배치한다.

    fit="cover"   영역을 꽉 채우고 넘치는 부분은 잘라낸다(사진용).
    fit="contain" 비율을 지켜 영역 안에 전부 들어가게 하고 남는 쪽은 비운다(로고용).
    """
    bx, by, bw, bh = box
    with Image.open(img_path) as im:
        iw, ih = im.size
    img_ratio = iw / ih
    box_ratio = bw / bh

    if fit == "contain":
        if img_ratio > box_ratio:
            w, h = bw, int(bw / img_ratio)
        else:
            h, w = bh, int(bh * img_ratio)
        x = bx + (bw - w) // 2
        y = by + (bh - h) // 2
        return slide.shapes.add_picture(str(img_path), Emu(x), Emu(y), width=Emu(w), height=Emu(h))

    # cover: 넘치는 부분은 crop 속성으로 잘라낸다(슬라이드 밖으로 밀어내지 않는다)
    pic = slide.shapes.add_picture(str(img_path), Emu(bx), Emu(by), width=Emu(bw), height=Emu(bh))
    if img_ratio > box_ratio:          # 이미지가 더 넓다 → 좌우를 자른다
        keep = box_ratio / img_ratio
        side = (1 - keep) / 2
        pic.crop_left = side
        pic.crop_right = side
    elif img_ratio < box_ratio:        # 이미지가 더 좁다 → 위아래를 자른다
        keep = img_ratio / box_ratio
        side = (1 - keep) / 2
        pic.crop_top = side
        pic.crop_bottom = side
    return pic


def grid_boxes(box, count, cols, gap_ratio=0.08):
    """box를 cols열 그리드로 나눠 count개의 셀 사각형을 만든다."""
    bx, by, bw, bh = box
    cols = max(1, cols)
    rows = (count + cols - 1) // cols
    gx = int(bw / cols * gap_ratio)
    gy = int(bh / rows * gap_ratio)
    cw = (bw - gx * (cols - 1)) // cols
    ch = (bh - gy * (rows - 1)) // rows
    cells = []
    for i in range(count):
        r, c = divmod(i, cols)
        cells.append((bx + c * (cw + gx), by + r * (ch + gy), cw, ch))
    return cells


def apply_overlay(prs, entries, base_dir: Path):
    """overlay 배치표대로 사용자 이미지를 각 슬라이드 위에 얹는다."""
    slides = list(prs.slides)
    placed = 0
    for entry in entries:
        n = entry.get("slide")
        if not isinstance(n, int) or not (1 <= n <= len(slides)):
            sys.exit(f"overlay: slide 번호가 범위를 벗어났습니다(1~{len(slides)}): {entry}")
        slide = slides[n - 1]
        fit = entry.get("fit", "cover")
        if fit not in ("cover", "contain"):
            sys.exit(f"overlay: fit은 cover 또는 contain이어야 합니다: {entry}")
        box = resolve_box(entry.get("box"))

        names = entry.get("images") or ([entry["image"]] if entry.get("image") else [])
        if not names:
            sys.exit(f"overlay: image 또는 images가 필요합니다: {entry}")
        paths = []
        for name in names:
            p = Path(name)
            if not p.is_absolute():
                p = base_dir / p
            if not p.exists():
                sys.exit(f"overlay: 이미지를 찾을 수 없습니다: {p}")
            paths.append(p)

        if len(paths) == 1:
            place_image(slide, paths[0], box, fit)
        else:
            cells = grid_boxes(box, len(paths), entry.get("cols", 3))
            for p, cell in zip(paths, cells):
                place_image(slide, p, cell, fit)
        placed += len(paths)
    return placed


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("images_dir", help="슬라이드 이미지가 들어 있는 폴더")
    ap.add_argument("-o", "--output", default="presentation.pptx")
    ap.add_argument("--notes", help="슬라이드별 발표자 노트 파일(--- 구분)")
    ap.add_argument("--overlay", help="사용자 첨부 이미지 배치표(JSON)")
    args = ap.parse_args()

    img_dir = Path(args.images_dir)
    images = sorted((p for p in img_dir.iterdir() if p.suffix.lower() in EXTS), key=natural_key)
    if not images:
        sys.exit(f"{img_dir} 안에 이미지 파일이 없습니다.")

    notes = load_notes(Path(args.notes)) if args.notes else []

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    full_box = (0, 0, int(SLIDE_W), int(SLIDE_H))
    for i, img_path in enumerate(images):
        slide = prs.slides.add_slide(blank)
        place_image(slide, img_path, full_box, "cover")
        if i < len(notes) and notes[i]:
            slide.notes_slide.notes_text_frame.text = notes[i]

    placed = 0
    if args.overlay:
        overlay_path = Path(args.overlay)
        entries = json.loads(overlay_path.read_text(encoding="utf-8"))
        if not isinstance(entries, list):
            sys.exit("overlay 파일은 배열(JSON list)이어야 합니다.")
        placed = apply_overlay(prs, entries, overlay_path.parent)

    prs.save(args.output)
    tail = f", 첨부 이미지 {placed}장" if placed else ""
    print(f"완료: {args.output} ({len(images)}장{tail})")


if __name__ == "__main__":
    main()
