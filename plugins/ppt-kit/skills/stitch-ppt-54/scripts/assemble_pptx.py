#!/usr/bin/env python3
"""Stitch에서 내보낸 슬라이드 이미지를 5:4 pptx로 조립한다.

사용법:
    python assemble_pptx.py <이미지폴더> -o 발표자료.pptx
    python assemble_pptx.py <이미지폴더> -o 발표자료.pptx --notes notes.txt

- 이미지 폴더 안의 png/jpg/jpeg/webp 파일을 자연 정렬 순서로 슬라이드에 넣는다.
- 슬라이드 크기는 5:4 (13.333 x 10.667 인치).
- 이미지 비율이 5:4와 다르면 여백 없이 중앙을 채우도록(cover) 크롭 배치한다.
- --notes 파일은 슬라이드별 발표자 노트. 슬라이드 사이를 '---' 한 줄로 구분한다.
"""
import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit("python-pptx가 없습니다. 먼저 실행: pip install python-pptx --break-system-packages")

try:
    from PIL import Image
except ImportError:
    sys.exit("Pillow가 없습니다. 먼저 실행: pip install Pillow --break-system-packages")

SLIDE_W = Emu(12192000)  # 13.333 inch
SLIDE_H = Emu(9753600)   # 10.667 inch (5:4)
EXTS = {".png", ".jpg", ".jpeg", ".webp"}


def natural_key(p: Path):
    return [int(t) if t.isdigit() else t.lower() for t in re.split(r"(\d+)", p.name)]


def load_notes(path: Path):
    text = path.read_text(encoding="utf-8")
    return [chunk.strip() for chunk in re.split(r"^\s*---\s*$", text, flags=re.M)]


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("images_dir", help="슬라이드 이미지가 들어 있는 폴더")
    ap.add_argument("-o", "--output", default="presentation.pptx")
    ap.add_argument("--notes", help="슬라이드별 발표자 노트 파일(--- 구분)")
    args = ap.parse_args()

    img_dir = Path(args.images_dir)
    images = sorted((p for p in img_dir.iterdir() if p.suffix.lower() in EXTS), key=natural_key)
    if not images:
        sys.exit(f"{img_dir} 안에 이미지 파일이 없습니다.")

    notes = load_notes(Path(args.notes)) if args.notes else []

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for i, img_path in enumerate(images):
        slide = prs.slides.add_slide(blank)
        with Image.open(img_path) as im:
            iw, ih = im.size
        slide_ratio = SLIDE_W / SLIDE_H
        img_ratio = iw / ih
        if abs(img_ratio - slide_ratio) < 0.01:
            slide.shapes.add_picture(str(img_path), 0, 0, width=SLIDE_W, height=SLIDE_H)
        elif img_ratio > slide_ratio:
            # 이미지가 더 넓다: 높이를 맞추고 좌우를 잘라낸다
            w = int(SLIDE_H * img_ratio)
            slide.shapes.add_picture(str(img_path), Emu(-(w - SLIDE_W) // 2), 0,
                                     width=Emu(w), height=SLIDE_H)
        else:
            # 이미지가 더 좁다: 너비를 맞추고 위아래를 잘라낸다
            h = int(SLIDE_W / img_ratio)
            slide.shapes.add_picture(str(img_path), 0, Emu(-(h - SLIDE_H) // 2),
                                     width=SLIDE_W, height=Emu(h))
        if i < len(notes) and notes[i]:
            slide.notes_slide.notes_text_frame.text = notes[i]

    prs.save(args.output)
    print(f"완료: {args.output} ({len(images)}장)")


if __name__ == "__main__":
    main()
